import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

# פונקציה להמיר עמודים ב-PDF לתמונות (ללא שמירת קבצים זמניים)
def pdf_to_images_in_memory(pdf_path):
    doc = fitz.open(pdf_path)
    images = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=300)  # שמירה על רזולוציה גבוהה
        images.append(pix)

    return images

# פונקציה ליצירת מצגת PowerPoint
def create_ppt_from_pixmaps(pixmaps, output_ppt_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # יחס 16:9
    prs.slide_height = Inches(7.5)

    for pix in pixmaps:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # יצירת שקופית ריקה
        image_stream = pix.tobytes("png")

        # שמירת התמונה בזיכרון כקובץ זמני
        with open("temp_image.png", "wb") as temp_image:
            temp_image.write(image_stream)

        img = slide.shapes.add_picture("temp_image.png", Inches(0), Inches(0), width=None, height=None)

        # התאמת גודל התמונה לגודל השקופית
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # חישוב יחסי והתאמה
        scale_w = slide_width / img.width
        scale_h = slide_height / img.height
        scale = min(scale_w, scale_h)

        img.width = int(img.width * scale)
        img.height = int(img.height * scale)
        img.left = int((slide_width - img.width) / 2)
        img.top = int((slide_height - img.height) / 2)

    prs.save(output_ppt_path)

# נתיב לקובץ ה-PDF
pdf_path = r"C:\Users\ASUS\Documents\code_images\overleaf_images\presentation\fina_presentation\presentation pdf.pdf"
output_ppt_path = r"C:\Users\ASUS\Documents\code_images\overleaf_images\presentation\fina_presentation\presentation pp.pptx"

# המרת PDF לתמונות בזיכרון
pixmaps = pdf_to_images_in_memory(pdf_path)

# יצירת מצגת PowerPoint
create_ppt_from_pixmaps(pixmaps, output_ppt_path)

print(f"מצגת נשמרה בהצלחה: {output_ppt_path}")
