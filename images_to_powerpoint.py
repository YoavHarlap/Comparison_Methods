import os
from pptx import Presentation
from pptx.util import Inches

# הגדרת תיקיית התמונות
image_folder = r"C:\Users\ASUS\Documents\code_images\overleaf_images\presentation\final_presentation\reflection_frames1"

# יצירת מצגת חדשה
prs = Presentation()

# הגדרת גודל השקופיות ל-16:9 (1920x1080)
prs.slide_width = Inches(13.33)  # 1920 / 144
prs.slide_height = Inches(7.5)  # 1080 / 144

# לולאה על כל הקבצים בתיקייה
for filename in sorted(os.listdir(image_folder)):
    if filename.lower().endswith(('.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tiff')):
        img_path = os.path.join(image_folder, filename)

        # יצירת שקופית חדשה עם פריסה ריקה
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # הוספת התמונה לשקופית (גודל מקורי)
        img = slide.shapes.add_picture(img_path, 0, 0)

        # המרת ערכי left ו-top לערכים שלמים (int)
        img.left = int((prs.slide_width - img.width) / 2)
        img.top = int((prs.slide_height - img.height) / 2)

# שמירת המצגת
pptx_output_path = os.path.join(image_folder, "presentation_16_9_centered.pptx")
prs.save(pptx_output_path)

print(f"מצגת PowerPoint נשמרה ב: {pptx_output_path}")
